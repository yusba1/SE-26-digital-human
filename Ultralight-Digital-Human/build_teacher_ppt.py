from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT_PPT = ROOT / "指导老师汇报_Ultralight数字人项目.pptx"
IMAGE_PATH = Path(
    "/Users/wangzhijie/.cursor/projects/Users-wangzhijie-Documents-Internship-Project-Ultralight-Digital-Human/assets/digital_human_training_streaming_diagram.png"
)


def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(34)
    title.text_frame.paragraphs[0].font.bold = True


def add_bullets(slide, lines, left=0.8, top=1.7, width=11.4, height=5.2, font_size=24):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
        p.space_after = Pt(10)


def add_speaker_note(slide, note_text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = note_text


def build():
    prs = Presentation()
    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_blank = prs.slide_layouts[6]

    # 1 封面
    slide = prs.slides.add_slide(layout_title)
    slide.shapes.title.text = "Ultralight Digital Human\n项目原理与工程实现"
    sub = slide.placeholders[1]
    sub.text = "轻量化口型驱动 · 三损失联合训练 · 流式推理\n汇报人：____    日期：____"
    add_speaker_note(slide, "老师好，我今天汇报的主题是Ultralight Digital Human项目。重点讲三块：模型原理、训练机制、流式部署。")

    # 2 背景与目标
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "1. 项目背景与目标")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• 背景：数字人项目常见问题是模型重、时延高、工程链路不完整",
            "• 目标：做一个可训练、可推理、可导出、可流式的轻量数字人系统",
            "• 输入：人物训练视频 + 测试音频；输出：音画同步说话视频",
            "• 价值：兼顾效果、实时性和工程可落地性",
        ],
    )
    add_speaker_note(slide, "本项目的定位不是单点模型，而是完整链路：从数据到部署到流式服务。")

    # 3 总体架构
    slide = prs.slides.add_slide(layout_blank)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = "2. 总体技术架构（端到端闭环）"
    p.font.bold = True
    p.font.size = Pt(34)
    if IMAGE_PATH.exists():
        slide.shapes.add_picture(str(IMAGE_PATH), Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.8))
    else:
        add_bullets(
            slide,
            [
                "• 预处理 -> SyncNet训练(可选) -> 主模型训练 -> 推理 -> ONNX导出 -> 流式服务",
                "• 一键入口：run_pipeline.sh",
            ],
            left=0.9,
            top=1.8,
            width=11.0,
            height=3.5,
            font_size=28,
        )
    add_speaker_note(slide, "这一页强调工程闭环。流水线把7个步骤串起来，适合复现和交付。")

    # 4 预处理
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "3. 数据与预处理约束")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• 输入质量要求：正脸稳定、低噪声、口齿清晰（质量直接决定上限）",
            "• 帧率约束：wenet 对应 20fps；hubert 对应 25fps",
            "• 预处理产物：full_body_img/、landmarks/、aud_wenet.npy 或 aud_hu.npy",
            "• 作用：把“原始视频+音频”变成可对齐、可训练的数据样本",
        ],
    )
    add_speaker_note(slide, "这里建议明确告诉老师：这个项目对数据质量非常敏感。")

    # 5 U-Net与6通道
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "4. 主模型：6通道 U-Net + 音频融合")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• 6通道来源：参考帧(3通道) + 当前遮挡帧(3通道) 按通道拼接",
            "• U-Net思想：先下采样“看全局”，再上采样“补细节”，并用跳连保真",
            "• 音频分支：提取音频语义后在瓶颈层与图像语义融合",
            "• 输出：3通道嘴部区域，再贴回原图实现说话效果",
        ],
    )
    add_speaker_note(slide, "解释U-Net时建议用“修图”类比：先看全局，再补局部。")

    # 6 结构可解释
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "5. 网络不是黑盒：结构与轻量化")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• 编码器：inc + down1~down4；解码器：up1~up4；输出层：outc",
            "• 通道规模（当前实现）：32 -> 64 -> 128 -> 256 -> 512",
            "• 轻量化模块：InvertedResidual / 深度可分离思路，降低计算量",
            "• 结论：通用U-Net架构 + 口型任务定制，不是新造名词",
        ],
    )
    add_speaker_note(slide, "老师若问创新点，可以答：创新在任务化融合与工程化链路，不在基础网络命名。")

    # 7 三损失
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "6. 三损失联合训练机制")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• 像素损失：L1(pred, gt)，保证嘴型位置与几何准确",
            "• 感知损失：MSE(VGG19(pred), VGG19(gt))，提升纹理与真实感",
            "• 同步损失：BCE(cos(a, v), 1)，约束音频与口型语义一致",
            "• 总损失：L = L_pixel + λp·L_perc + λs·L_sync",
        ],
    )
    add_speaker_note(slide, "可强调三者分工：位置、质感、同步。")

    # 8 感知损失解释
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "7. 感知损失 vs 像素损失")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• VGG19在这里是“固定特征提取器”，不是分类器",
            "• 特征含义：边缘、纹理、局部结构等高层视觉表征",
            "• 像素损失偏“数值一致”，感知损失偏“观感一致”",
            "• 两者互补：避免只对齐位置却画面发糊",
        ],
    )
    add_speaker_note(slide, "如果老师追问特征空间，可以说：比较的是中间层特征图而非RGB逐像素。")

    # 9 同步损失解释
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "8. 同步网络与余弦相似度")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• SyncNet含两分支：face_encoder(口型图像) 与 audio_encoder(音频特征)",
            "• 两分支输出同维embedding，归一化后计算 cos(a, v)",
            "• 关键点：不是原始音频和图像比较，而是语义向量比较",
            "• 用 BCE(cos(a,v),1) 训练，使匹配音画相似度持续提升",
        ],
    )
    add_speaker_note(slide, "强调‘同维是可比较前提；损失函数才让相似度变大’。")

    # 10 两个onnx
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "9. encoder.onnx 与 unet.onnx 的职责")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• encoder.onnx：实时音频编码器（提取声学语义特征）",
            "• unet.onnx：口型生成器（图像上下文 + 音频特征 -> 嘴部像素）",
            "• 串联关系：先编码音频，再触发口型生成",
            "• 这样拆分有利于实时服务和跨端部署",
        ],
    )
    add_speaker_note(slide, "一言以蔽之：encoder负责“听懂”，unet负责“画出来”。")

    # 11 流式时序
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "10. 流式推理时序：10ms输入，窗口触发出图")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• 10ms定义：16kHz下每160个采样点",
            "• 不是每10ms都出图：10ms是输入粒度，不是渲染粒度",
            "• 音频累计达到阈值后触发一次encoder；特征窗口满足后触发一次unet",
            "• 核心目标：在实时性、稳定性、同步性之间做工程平衡",
        ],
    )
    add_speaker_note(slide, "这页用于解释‘为什么看起来连续，但内部是事件触发式生成’。")

    # 12 结果与展望
    slide = prs.slides.add_slide(layout_content)
    set_title(slide, "11. 结果、局限与下一步")
    slide.placeholders[1].text = ""
    add_bullets(
        slide,
        [
            "• 已完成：一键流水线、可训练、可推理、可ONNX、可流式",
            "• 局限：对数据质量敏感；同步损失可继续引入更强对比学习",
            "• 下一步：模型压缩量化、端侧优化、统一评估指标体系",
            "• 目标：从可用原型走向稳定工程化交付",
        ],
    )
    add_speaker_note(slide, "收尾建议强调：该项目已具备工程转化基础，下一步是性能与稳健性迭代。")

    # 13 Q&A
    slide = prs.slides.add_slide(layout_blank)
    t = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.0), Inches(2.2))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Q & A"
    p.font.size = Pt(72)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_speaker_note(slide, "谢谢老师，欢迎提问。")

    prs.save(str(OUTPUT_PPT))
    print(f"已生成: {OUTPUT_PPT}")


if __name__ == "__main__":
    build()
